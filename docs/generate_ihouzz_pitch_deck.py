# -*- coding: utf-8 -*-
"""Generate MindX ITBA Level 3 pitch deck: iHouzz centralized warehouse + listing."""
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
COBALT = RGBColor(0, 71, 171)
COBALT_LIGHT = RGBColor(230, 240, 255)
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(33, 37, 41)

FONT_VI = "Segoe UI"
FONT_TITLE = "Segoe UI"


def make_logo(path: Path) -> None:
    w, h = 520, 120
    im = Image.new("RGB", (w, h), (0, 71, 171))
    draw = ImageDraw.Draw(im)
    try:
        font = ImageFont.truetype("C:/Windows/Fonts/segoeui.ttf", 44)
        font_sm = ImageFont.truetype("C:/Windows/Fonts/segoeui.ttf", 18)
    except OSError:
        font = ImageFont.load_default()
        font_sm = font
    draw.text((28, 28), "iHouzz.com", fill=(255, 255, 255), font=font)
    draw.text((28, 86), "Nội bộ · Kho & niêm yết", fill=(200, 220, 255), font=font_sm)
    im.save(path)


def set_run_font(run, size_pt: int, bold: bool = False, color=TEXT_DARK):
    run.font.name = FONT_VI
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_only_slide(prs, title: str, subtitle: str | None, dark_bg: bool):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    if dark_bg:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COBALT
    left, top, width, height = Inches(0.6), Inches(1.9), Inches(12.3), Inches(1.4)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_run_font(p.runs[0], 36, True, WHITE if dark_bg else COBALT)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.space_before = Pt(12)
        set_run_font(p2.runs[0], 20, False, RGBColor(220, 230, 255) if dark_bg else TEXT_DARK)
    return slide


def add_bullet_slide(prs, title: str, bullets: list[str], notes: str):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COBALT_LIGHT

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.25),
        Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COBALT
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.5), Inches(0.85))
    tfp = tb.text_frame
    tfp.paragraphs[0].text = title
    tfp.paragraphs[0].alignment = PP_ALIGN.LEFT
    set_run_font(tfp.paragraphs[0].runs[0], 28, True, COBALT)

    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(12.4), Inches(5.6))
    bf = body.text_frame
    bf.word_wrap = True
    bf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = line
        para.level = 0
        para.space_after = Pt(10)
        para.line_spacing = 1.15
        set_run_font(para.runs[0], 18, False, TEXT_DARK)
        para.font.name = FONT_VI

    if notes:
        ns = slide.notes_slide
        ns.notes_text_frame.text = notes
    return slide


def add_logo_header(slide, logo_path: Path):
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(11.95), Inches(0.35), height=Inches(0.55))


def main():
    logo_path = OUT_DIR / "ihouzz_logo_pitch.png"
    make_logo(logo_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_spec = [
        (
            "cover",
            {
                "title": "Kho tập trung & niêm yết\nmột nền tảng cho tốc độ iHouzz",
                "subtitle": "Thiết kế quy trình — MindX ITBA Level 3 · Đối tác doanh nghiệp iHouzz",
            },
        ),
        (
            "agenda",
            {
                "title": "Hôm nay chúng ta đi theo bốn nhịp",
                "bullets": [
                    "Mở đầu: bối cảnh và câu hỏi kinh doanh cần giải.",
                    "Thực trạng: pain theo từng phòng — MKT, KD, KD hệ thống, IT.",
                    "Giải pháp: quy trình 2 tầng, phân quyền, kiểm toán, demo ánh xạ BR.",
                    "Kết luận: hiệu quả, lộ trình, và cách chúng ta cùng làm tiếp.",
                ],
                "notes": "Gợi ý hình: sơ đồ timeline ngang 4 mốc với icon nhẹ (loupe, nhóm người, gear, flag).",
            },
        ),
        (
            "intro",
            {
                "title": "Góc nhìn ITBA: cầu nối giữa nghiệp vụ và sản phẩm",
                "bullets": [
                    "Chúng tôi dịch yêu cầu iHouzz thành luồng trạng thái rõ ràng, có vai phê duyệt.",
                    "Mỗi pain phòng ban được gắn mã BR/SRS để IT triển khai không mơ hồ.",
                    "Prototype React + API mock minh họa đúng thứ tự thao tác người dùng thật.",
                ],
                "notes": "Gợi ý hình: silhouette BA giữa hai khối “Business” và “Engineering” nối bằng mũi tên.",
            },
        ),
        (
            "ask",
            {
                "title": "Doanh nghiệp cần gì — trong một câu",
                "bullets": [
                    "Một kho BĐS nội bộ thống nhất, dễ tìm — dễ chọn — để đẩy tin lên sàn iHouzz.",
                    "Quy trình giảm nhập sai, rút ngắn vòng chờ duyệt, giữ uy tín nội dung public.",
                    "Định hướng tính năng: không chỉ màn hình, mà là quy tắc vận hành đo được.",
                ],
                "notes": "Gợi ý hình: funnel từ “Nguồn hàng” xuống “Kho chuẩn” rồi “Niêm yết sàn”.",
            },
        ),
        (
            "context",
            {
                "title": "Thực trạng: một sản phẩm BĐS, nhiều bàn tay",
                "bullets": [
                    "Đầu chủ/KD nhập thông tin; GĐ POS xác nhận chất lượng kho; MKT kiểm duyệt tin.",
                    "IT vận hành dữ liệu, tích hợp sàn; bộ phận phát triển KD định nghĩa luật kinh doanh.",
                    "Nếu thiếu khung trạng thái, cùng một căn có thể bị hiểu khác nhau giữa các phòng.",
                ],
                "notes": "Gợi ý hình: swimlane đơn giản 4 hàng (Sales, POS, MKT, IT) với một “thẻ tài sản” đi qua.",
            },
        ),
        (
            "pain_mkt",
            {
                "title": "Pain Marketing: uy tín sàn nằm ở lớp kiểm duyệt",
                "bullets": [
                    "Cần queue tin rõ ràng, lịch sử sửa, và tiêu chí từ chối minh bạch cho đối tác nội bộ.",
                    "Niêm yết quá chậm làm mất cơ hội rao bán; niêm yết ẩu là rủi ro pháp lý và thương hiệu.",
                    "Thiếu liên kết với trạng thái kho thì MKT khó biết tin đang “đứng trên nền” nào.",
                ],
                "notes": "Gợi ý hình: biểu tượng khiên + checklist 3 ô (hình ảnh, giá, pháp lý) màu coban.",
            },
        ),
        (
            "pain_sales",
            {
                "title": "Pain Kinh doanh: tốc độ không được mua bằng độ chính xác",
                "bullets": [
                    "Trùng địa chỉ, trùng nguồn gây tranh chấp hoa hồng và mất niềm tin nội bộ.",
                    "Đầu chủ cần biết “đang ở bước nào” thay vì hỏi lại trợ lý qua chat nhiều lần.",
                    "Gỡ tin / gỡ nguồn nếu không có luật thì dễ để lại “tin ma” hoặc dữ liệu không khớp sàn.",
                ],
                "notes": "Gợi ý hình: đồng hồ cạnh dấu tick xanh — nhấn mạnh cân bằng tốc độ và đúng.",
            },
        ),
        (
            "pain_bizdev",
            {
                "title": "Pain phát triển hệ thống KD: luật nghiệp vụ phải viết được",
                "bullets": [
                    "Cascade gỡ nguồn, chặn thao tác khi đang niêm yết — cần mô tả bằng BR, không chỉ ý miệng.",
                    "Đồng bộ Level 1 / Level 2 giữa kho và tin là điểm dễ lỗi nếu không có sơ đồ trạng thái.",
                    "Ma trận phòng ban giúp đội KD hệ thống thống nhất “được / không được” với sàn điện tử.",
                ],
                "notes": "Gợi ý hình: state machine đơn giản 4–5 node (nháp) tô màu coban.",
            },
        ),
        (
            "pain_it",
            {
                "title": "Pain IT: bảo mật, tải, và khả năng kiểm toán sau này",
                "bullets": [
                    "Phân quyền theo POS: nhân viên không được nhìn thẳng địa chỉ chi nhánh không thuộc phạm vi.",
                    "Payload lớn, ghi log nhiều — cần thiết kế API và audit không làm nghẽn vận hành.",
                    "No hard delete: giữ vết để phục vụ tranh chấp và cứu hộ dữ liệu an toàn hơn xóa vĩnh viễn.",
                ],
                "notes": "Gợi ý hình: icon khóa trên bản đồ Việt Nam mờ (masking) vài khu vực.",
            },
        ),
        (
            "impact",
            {
                "title": "Hệ quả vận hành khi chưa có khung rõ",
                "bullets": [
                    "Sai sót nhập liệu và trùng nguồn làm chậm chu kỳ từ kho đến tin public.",
                    "Trải nghiệm nội bộ kém: nhiều hệ, nhiều hỏi đáp, khó đo SLA phòng ban.",
                    "Rủi ro uy tín sàn: tin không khớp thực tế hoặc gỡ không sạch khỏi kênh điện tử.",
                ],
                "notes": "Gợi ý hình: biểu đồ cột “thời gian chờ duyệt” cao vs “tỉ lệ lỗi nhập” cao — hai cột đỏ nhạt.",
            },
        ),
        (
            "gap",
            {
                "title": "Khoảng trống: thiếu một “bản đồ quy trình” chung",
                "bullets": [
                    "Cần một lớp trung gian: trạng thái kho và trạng thái niêm yết song song nhưng đồng bộ.",
                    "Cần ranh giới phê duyệt: ai được tạo, ai được đẩy tin, ai được gỡ, ai được xem dữ liệu gốc.",
                    "Cần bằng chứng định lượng: log hành động, export phục vụ kiểm tra sau sự cố.",
                ],
                "notes": "Gợi ý hình: khoảng trống giữa hai khối puzzle — slide tiếp theo lắp vào.",
            },
        ),
        (
            "sol_two_level",
            {
                "title": "Giải pháp trụ cột 1: kiến trúc hai tầng trạng thái",
                "bullets": [
                    "Tầng kho (Level 1) mô tả vòng đời tài sản nội bộ; tầng tin (Level 2) mô tả hiển thị sàn.",
                    "Auto-sync có điều kiện: thay đổi tin/kho phải tuân quy tắc đã chốt trong BRD.",
                    "Demo F3, F5, F7, F8 trên prototype cho thấy thứ tự thao tác khớp swimlane nghiệp vụ.",
                ],
                "notes": "Gợi ý hình: hai vòng tròn lồng hoặc hai thanh timeline song song (Kho | Tin).",
            },
        ),
        (
            "sol_rbac",
            {
                "title": "Giải pháp trụ cột 2: RBAC động & che dữ liệu chi nhánh",
                "bullets": [
                    "Admin cấu hình ma trận quyền; nhân viên chỉ thấy thao tác và cột dữ liệu được phép.",
                    "Data masking địa chỉ theo POS giảm lộ thông tin giữa chi nhánh cạnh tranh nội bộ.",
                    "Demo F9–F10: cùng màn hình nhưng khác dữ liệu hiển thị theo role — dễ chứng minh cho IT.",
                ],
                "notes": "Gợi ý hình: bảng ma trận Role × Permission với ô tick màu coban.",
            },
        ),
        (
            "sol_audit",
            {
                "title": "Giải pháp trụ cột 3: audit trail & không xóa cứng",
                "bullets": [
                    "Mỗi phê duyệt, từ chối, gỡ tin, gỡ nguồn ghi log bất biến, có thể export CSV.",
                    "Giám khảo và iHouzz có thể truy vết “ai làm gì lúc nào” phục vụ họp rà soát.",
                    "Demo F11: lọc theo thời gian và actor — đúng tư duy kiểm soát nội bộ enterprise.",
                ],
                "notes": "Gợi ý hình: cuộn giấy cổ điển hoặc icon journal với dòng thời gian dọc.",
            },
        ),
        (
            "flow_f2",
            {
                "title": "Luồng vào kho: nhập liệu thông minh, giảm trùng sớm",
                "bullets": [
                    "SmartAddress + debounce kiểm tra trùng địa chỉ trước khi tạo bản ghi kho.",
                    "Form chuẩn hóa loại giao dịch, diện tích, pháp lý — giảm rework cho MKT sau này.",
                    "Liên kết ITBA: acceptance criteria rõ cho API search trùng và ngưỡng cảnh báo.",
                ],
                "notes": "Gợi ý hình: mock UI ô địa chỉ với badge “không trùng” / “cảnh báo trùng”.",
            },
        ),
        (
            "flow_f3",
            {
                "title": "Luồng chất lượng kho: GĐ POS phân loại kho chuẩn / đảm bảo",
                "bullets": [
                    "Tách quyền: chỉ vai POS manager quyết định đưa tài sản vào kho niêm yết được hay chưa.",
                    "Ghi log và thông báo nội bộ sau phê duyệt — minh bạch với KD và MKT.",
                    "Demo F3: chọn tài sản, duyệt, quan sát cập nhật trạng thái kho ngay trên danh sách.",
                ],
                "notes": "Gợi ý hình: stamp “Approved” màu coban trên thẻ property card.",
            },
        ),
        (
            "flow_f4f5",
            {
                "title": "Luồng lên sàn: soạn tin — hàng chờ — trung tâm duyệt niêm yết",
                "bullets": [
                    "Sales soạn nội dung public; MKT/Admin duyệt để bảo vệ chất lượng sàn iHouzz.",
                    "Trạng thái tin rõ ràng: nháp, chờ duyệt, đang niêm yết — tránh đăng nhầm sớm.",
                    "Demo F4 → F5: một mạch từ gửi duyệt đến phê duyệt, phản ánh UC005 trong SRS.",
                ],
                "notes": "Gợi ý hình: pipeline ngang 4 bước với icon megaphone và patch-check.",
            },
        ),
        (
            "flow_f6f7",
            {
                "title": "Luồng gỡ tin: yêu cầu có kiểm soát, đồng bộ Level 2",
                "bullets": [
                    "Sales khởi tạo yêu cầu gỡ; Marketing phê duyệt — tránh gỡ một chiều không giải thích.",
                    "Sau duyệt, hệ thống cập nhật trạng thái tin và kho theo quy tắc auto-sync đã định nghĩa.",
                    "Demo F6–F7: thấy rõ người gửi, người duyệt, thời điểm — phục vụ họp MKT–KD.",
                ],
                "notes": "Gợi ý hình: hai nút bấm “Request” và “Approve” nối bằng mũi tên có nhãn thời gian.",
            },
        ),
        (
            "flow_f8",
            {
                "title": "Luồng gỡ nguồn: luật cascade & chặn khi đang niêm yết",
                "bullets": [
                    "BR-010: không cho gỡ nguồn khi tin còn trên sàn — giảm rủi ro dữ liệu “mồ côi”.",
                    "Khi được phép, cascade hủy tin liên quan — ITBA diễn giải như một use case có điều kiện.",
                    "Demo F8: thử case bị chặn và case được duyệt để giám khảo thấy luật chạy đúng.",
                ],
                "notes": "Gợi ý hình: sơ đồ if/else (diamond) màu coban: “Đang niêm yết?” Yes → stop.",
            },
        ),
        (
            "flow_f9",
            {
                "title": "Giám sát kho đa POS: đúng phạm vi, đúng độ nhạy dữ liệu",
                "bullets": [
                    "Dashboard kho tổng hợp nhưng ẩn địa chỉ chi nhánh khác theo BR-013.",
                    "Giúp lãnh đạo xem sức khỏe kho mà không vi phạm ranh giới tổ chức.",
                    "Demo F9: đổi role hoặc POS trong login để thấy masking đổi theo — bằng chứng trực quan.",
                ],
                "notes": "Gợi ý hình: bản đồ pin bị làm mờ (blur) một nửa địa chỉ.",
            },
        ),
        (
            "flow_f10",
            {
                "title": "IAM & POS: ma trận quyền là “hợp đồng” giữa KD hệ thống và IT",
                "bullets": [
                    "User gắn POS, trạng thái hoạt động; quyền chi tiết đến từng hành động trên UI.",
                    "Thay đổi ma trận phản ánh ngay menu — giảm phụ thuộc release cứng cho từng nhân sự.",
                    "Demo F10: ví dụ tắt quyền tạo tài sản cho một sales test — tái hiện change request thật.",
                ],
                "notes": "Gợi ý hình: icon people-network với các nhánh quyền màu coban.",
            },
        ),
        (
            "demo",
            {
                "title": "Minh chứng: prototype đã ánh xạ BRD / SRS / ma trận quyền",
                "bullets": [
                    "React + json-server: mô phỏng contract API để IT review sớm trước backend production.",
                    "Các màn F2–F12 khớp thứ tự kịch bản UAT — giám khảo có thể tái hiện sau buổi trình bày.",
                    "Đăng nhập demo: ví dụ admin@ihouzz.com, mật khẩu 123456, OTP 111111 (môi trường học).",
                ],
                "notes": "Gợi ý hình: screenshot collage 3 màn (Login, F5, F11) trên nền trắng coban border.",
            },
        ),
        (
            "kpi",
            {
                "title": "Hiệu quả mong đợi: gần với chỉ số phòng ban quan tâm",
                "bullets": [
                    "Giảm lỗi nhập & trùng nguồn nhờ kiểm tra sớm và quy trình phê duyệt có log.",
                    "Rút thời gian từ kho chuẩn đến tin được duyệt nhờ queue rõ và trách nhiệm vai.",
                    "Trải nghiệm nội bộ tốt hơn: ít hỏi lại, SLA minh bạch, dữ liệu đồng bộ với sàn.",
                ],
                "notes": "Gợi ý hình: ba chỉ số KPI card (Error rate ↓, Time-to-list ↓, eNPS ↑) màu coban.",
            },
        ),
        (
            "roadmap",
            {
                "title": "Lộ trình gợi ý: ưu tiên theo rủi ro nghiệp vụ & nỗ lực IT",
                "bullets": [
                    "Giai đoạn 1: khóa state machine, API kho/tin, auth thật — nền cho mọi phòng.",
                    "Giai đoạn 2: tích hợp sàn thật, thông báo, báo cáo SLA cho MKT và KD.",
                    "Giai đoạn 3: tối ưu hiệu năng, retention log, và policy lưu trữ theo chuẩn công ty.",
                ],
                "notes": "Gợi ý hình: roadmap 3 quý dạng Gantt đơn giản với mốc “Go-live pilot POS”.",
            },
        ),
        (
            "closing",
            {
                "title": "Cảm ơn MindX & iHouzz — sẵn sàng đào sâu từng pain",
                "bullets": [
                    "Chúng tôi mang đến bản thiết kế quy trình có thể đo, có thể demo, có thể hand-off IT.",
                    "Mời giám khảo và doanh nghiệp đặt câu hỏi theo phòng: MKT, KD, KD hệ thống, hoặc IT.",
                    "Liên hệ: [tên nhóm] — prototype chạy local theo README dự án ihouzz-demo.",
                ],
                "notes": "Gợi ý hình: logo iHouzz + QR placeholder tới repo/demo (nếu bạn công khai).",
            },
        ),
    ]

    # Slide 1 cover
    s0 = add_title_only_slide(
        prs,
        slides_spec[0][1]["title"],
        slides_spec[0][1]["subtitle"],
        dark_bg=True,
    )
    s0.shapes.add_picture(str(logo_path), Inches(0.55), Inches(0.45), height=Inches(0.65))
    s0.notes_slide.notes_text_frame.text = (
        "Gợi ý hình: nền coban full-bleed, logo góc trái; có thể thêm pattern sóng nhẹ mờ trắng 5% opacity."
    )

    for kind, data in slides_spec[1:]:
        if kind == "cover":
            continue
        slide = add_bullet_slide(prs, data["title"], data["bullets"], data.get("notes", ""))
        add_logo_header(slide, logo_path)

    out_pptx = OUT_DIR / "iHouzz_Kho_NiemYet_MindX_ITBA_L3.pptx"
    prs.save(out_pptx)
    print("Saved:", out_pptx)


if __name__ == "__main__":
    main()
